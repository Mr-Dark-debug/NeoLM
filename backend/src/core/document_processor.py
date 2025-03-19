import os
import asyncio
import mimetypes
from typing import List, Dict, Any
from docling.document_converter import DocumentConverter
from moviepy.editor import VideoFileClip
import aiohttp
from crawl4ai import AsyncWebCrawler
from crawl4ai.async_configs import BrowserConfig, CrawlerRunConfig
from urllib.parse import urlparse
from src.core.config import Config
from src.utils.logging import logger
from pydantic import BaseModel

class ProcessingResult(BaseModel):
    success: bool
    content: str
    metadata: Dict[str, Any] = {}
    error_message: str = ""

class DocumentProcessor:
    def __init__(self):
        Config.validate()
        self.headers = {"Authorization": f"Bearer {Config.HF_API_KEY}"}
        self.converter = DocumentConverter()
        self.browser_config = BrowserConfig(verbose=False)
        self.crawler_run_config = CrawlerRunConfig(
            word_count_threshold=50,
            remove_overlay_elements=True,
            process_iframes=True
        )
        os.makedirs(Config.TMP_DIR, exist_ok=True)

    async def process_files(self, file_paths: List[str]) -> List[ProcessingResult]:
        tasks = [self.process_file(path) for path in file_paths]
        return await asyncio.gather(*tasks)

    def _is_supported_file_type(self, mime_type: str) -> bool:
        supported_types = {
            "text/plain", "text/csv", "text/markdown",
            "application/pdf",
            "application/msword",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-excel",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "audio/mpeg", "audio/wav",
            "video/mp4", "video/x-msvideo",
            "image/jpeg", "image/png", "image/gif"
        }
        return mime_type in supported_types or mime_type.startswith(("audio/", "video/", "image/"))

    async def process_file(self, file_path: str) -> ProcessingResult:
        try:
            if urlparse(file_path).scheme in ("http", "https"):
                return await self._process_url(file_path)
            
            if not os.path.exists(file_path):
                return ProcessingResult(success=False, content="", error_message=f"File not found: {file_path}")
            
            if not self._check_file_size(file_path):
                return ProcessingResult(success=False, content="", error_message=f"File {file_path} exceeds size limit")
            
            mime_type = self._get_mime_type(file_path)
            if not mime_type:
                return ProcessingResult(success=False, content="", error_message=f"Unknown file type: {file_path}")
            
            if not self._is_supported_file_type(mime_type):
                return ProcessingResult(success=False, content="", error_message=f"Unsupported file type: {mime_type}")
            
            logger.info(f"Processing {file_path} ({mime_type})")
            if mime_type.startswith("audio/"):
                return await self._process_audio(file_path)
            elif mime_type.startswith("video/"):
                return await self._process_video(file_path)
            elif mime_type.startswith("image/"):
                return await self._process_image(file_path)
            elif mime_type == "text/plain":
                # Handle plain text directly
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                return ProcessingResult(
                    success=True,
                    content=content,
                    metadata={"type": "text", "path": file_path}
                )
            else:
                return self._process_document(file_path)
        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}")
            return ProcessingResult(success=False, content="", error_message=str(e))

    def _process_document(self, file_path: str) -> ProcessingResult:
        try:
            mime_type = self._get_mime_type(file_path)
            if mime_type in ("application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"):
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    content = "\n".join(text)
                    return ProcessingResult(
                        success=True,
                        content=content,
                        metadata={"type": "presentation", "path": file_path}
                    )
                except ImportError:
                    return ProcessingResult(success=False, content="", error_message="PPTX processing requires python-pptx library")
            
            result = self.converter.convert(file_path)
            return ProcessingResult(
                success=True,
                content=result.text,
                metadata={"type": "document", "path": file_path}
            )
        except Exception as e:
            return ProcessingResult(success=False, content="", error_message=f"Docling error: {str(e)}")

    async def _process_url(self, file_path: str) -> ProcessingResult:
        try:
            async with AsyncWebCrawler(config=self.browser_config) as crawler:
                result = await crawler.arun(url=file_path, config=self.crawler_run_config)
                if result.success:
                    return ProcessingResult(
                        success=True,
                        content=result.markdown,
                        metadata={"type": "url", "path": file_path}
                    )
                raise RuntimeError("Failed to crawl URL")
        except Exception as e:
            return ProcessingResult(success=False, content="", error_message=str(e))

    async def _process_video(self, file_path: str) -> ProcessingResult:
        try:
            video = VideoFileClip(file_path)
            if video.duration > Config.MAX_AUDIO_LENGTH:
                video.close()
                return ProcessingResult(success=False, content="", error_message=f"Video exceeds {Config.MAX_AUDIO_LENGTH}s")
            
            temp_audio = os.path.join(Config.TMP_DIR, f"temp_{os.path.basename(file_path)}.wav")
            video.audio.write_audiofile(temp_audio, logger=None)
            result = await self._process_audio(temp_audio)
            video.close()
            os.unlink(temp_audio)
            return ProcessingResult(
                success=result.success,
                content=f"Video transcription:\n{result.content}",
                metadata={"type": "video", "path": file_path},
                error_message=result.error_message
            )
        except Exception as e:
            return ProcessingResult(success=False, content="", error_message=str(e))

    async def _process_audio(self, file_path: str) -> ProcessingResult:
        API_URL = "https://api-inference.huggingface.co/models/openai/whisper-large-v3-turbo"
        try:
            if os.path.getsize(file_path) > Config.MAX_FILE_SIZE:
                return ProcessingResult(success=False, content="", error_message="Audio file too large")
            
            with open(file_path, "rb") as f:
                async with aiohttp.ClientSession() as session:
                    async with session.post(API_URL, headers=self.headers, data=f) as resp:
                        result = await resp.json()
                        transcription = result.get("text", "")
                        if not transcription:
                            raise ValueError("No transcription received")
                        return ProcessingResult(
                            success=True,
                            content=transcription,
                            metadata={"type": "audio", "path": file_path}
                        )
        except Exception as e:
            return ProcessingResult(success=False, content="", error_message=str(e))

    async def _process_image(self, file_path: str) -> ProcessingResult:
        API_URL = "https://api-inference.huggingface.co/models/Salesforce/blip-image-captioning-large"
        try:
            with open(file_path, "rb") as f:
                async with aiohttp.ClientSession() as session:
                    async with session.post(API_URL, headers=self.headers, data=f) as resp:
                        result = await resp.json()
                        caption = result[0]["generated_text"] if isinstance(result, list) else result.get("error", "Image processing failed")
                        return ProcessingResult(
                            success=True,
                            content=caption,
                            metadata={"type": "image", "path": file_path}
                        )
        except Exception as e:
            return ProcessingResult(success=False, content="", error_message=f"Image processing error: {str(e)}")

    def _check_file_size(self, file_path: str) -> bool:
        if not os.path.exists(file_path):
            return False
        return os.path.getsize(file_path) <= Config.MAX_FILE_SIZE

    def _get_mime_type(self, file_path: str) -> str:
        mime_type = mimetypes.guess_type(file_path)[0]
        if mime_type is None:
            ext = os.path.splitext(file_path)[1].lower()
            mime_map = {
                ".txt": "text/plain",
                ".md": "text/markdown",
                ".mp3": "audio/mpeg",
                ".wav": "audio/wav",
                ".mp4": "video/mp4",
                ".avi": "video/x-msvideo",
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".png": "image/png",
                ".gif": "image/gif",
                ".pdf": "application/pdf",
                ".csv": "text/csv",
                ".doc": "application/msword",
                ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ".ppt": "application/vnd.ms-powerpoint",
                ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ".xls": "application/vnd.ms-excel",
                ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            }
            mime_type = mime_map.get(ext)
        return mime_type